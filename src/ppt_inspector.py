#!/usr/bin/env python3

import os
import json
import csv
import hashlib
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, asdict

import click
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.shapes.group import GroupShape
try:
    from pptx.shapes.autoshape import AutoShape
    from pptx.shapes.freeform import FreeformBuilder
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import MSO_AUTO_SIZE
except ImportError:
    # Fallback for older python-pptx versions
    MSO_SHAPE_TYPE = None
    AutoShape = None
    FreeformBuilder = None
    MSO_AUTO_SIZE = None
from PIL import Image


@dataclass
class SlideInfo:
    slide_number: int
    title: Optional[str]
    text_content: List[str]
    shape_count: int
    image_count: int
    image_files: List[str]
    text_shapes: List[Dict[str, Any]]
    graphic_elements: List[Dict[str, Any]]
    logos_and_brands: List[Dict[str, Any]]
    tags: List[str]
    notes: Optional[str]


@dataclass
class DocumentMetadata:
    filename: str
    file_size: int
    created_date: Optional[str]
    modified_date: Optional[str]
    author: Optional[str]
    title: Optional[str]
    subject: Optional[str]
    category: Optional[str]
    comments: Optional[str]
    slide_count: int
    total_images: int
    company_mentions: List[str]
    copyright_notices: List[str]
    confidentiality_labels: List[str]
    custom_tags: List[str]


class PowerPointInspector:
    def __init__(self, filepath: str, output_dir: str = "exports", image_dir: str = "images"):
        self.filepath = Path(filepath)
        self.output_dir = Path(output_dir)
        self.image_dir = Path(image_dir)
        self.presentation = None
        self.metadata = None
        self.slides_info = []

        self.output_dir.mkdir(exist_ok=True)
        self.image_dir.mkdir(exist_ok=True)

        # Common patterns for detection
        self.copyright_patterns = [
            "©", "copyright", "(c)", "©️", "all rights reserved"
        ]
        self.confidentiality_patterns = [
            "confidential", "proprietary", "internal", "restricted", "private",
            "not for distribution", "do not distribute", "internal use only"
        ]
        self.company_patterns = [
            "inc.", "llc", "corp.", "corporation", "ltd.", "limited", "company"
        ]

    def load_presentation(self) -> None:
        """Load PowerPoint presentation."""
        if not self.filepath.exists():
            raise FileNotFoundError(f"File not found: {self.filepath}")

        self.presentation = Presentation(str(self.filepath))
        print(f"Loaded presentation: {self.filepath.name}")
        print(f"Total slides: {len(self.presentation.slides)}")

    def extract_document_metadata(self) -> DocumentMetadata:
        """Extract document-level metadata and properties."""
        core_props = self.presentation.core_properties
        file_stats = self.filepath.stat()

        # Collect all text for pattern analysis
        all_text = self._get_all_text_content()

        metadata = DocumentMetadata(
            filename=self.filepath.name,
            file_size=file_stats.st_size,
            created_date=core_props.created.isoformat() if core_props.created else None,
            modified_date=core_props.modified.isoformat() if core_props.modified else None,
            author=core_props.author,
            title=core_props.title,
            subject=core_props.subject,
            category=core_props.category,
            comments=core_props.comments,
            slide_count=len(self.presentation.slides),
            total_images=self._count_total_images(),
            company_mentions=self._find_pattern_matches(all_text, self.company_patterns),
            copyright_notices=self._find_pattern_matches(all_text, self.copyright_patterns),
            confidentiality_labels=self._find_pattern_matches(all_text, self.confidentiality_patterns),
            custom_tags=[]
        )

        self.metadata = metadata
        return metadata

    def extract_slide_content(self) -> List[SlideInfo]:
        """Extract content from each slide."""
        slides_info = []

        for i, slide in enumerate(self.presentation.slides):
            slide_info = self._analyze_slide(slide, i + 1)
            slides_info.append(slide_info)

        self.slides_info = slides_info
        return slides_info

    def _analyze_slide(self, slide, slide_number: int) -> SlideInfo:
        """Analyze individual slide content."""
        text_content = []
        image_files = []
        text_shapes = []
        graphic_elements = []
        logos_and_brands = []
        shape_count = 0
        image_count = 0

        # Extract slide title
        title = self._extract_slide_title(slide)

        # Process all shapes
        for shape in slide.shapes:
            shape_count += 1
            shape_analysis = self._analyze_shape_comprehensive(shape, slide_number)

            # Collect text content
            if shape_analysis['text']:
                if isinstance(shape_analysis['text'], list):
                    text_content.extend(shape_analysis['text'])
                else:
                    text_content.append(shape_analysis['text'])

            # Collect images
            if shape_analysis['image_file']:
                if isinstance(shape_analysis['image_file'], list):
                    image_files.extend(shape_analysis['image_file'])
                    image_count += len(shape_analysis['image_file'])
                else:
                    image_files.append(shape_analysis['image_file'])
                    image_count += 1

            # Collect text shapes (logos, titles, formatted text)
            if shape_analysis['text_shape']:
                if isinstance(shape_analysis['text_shape'], list):
                    text_shapes.extend(shape_analysis['text_shape'])
                else:
                    text_shapes.append(shape_analysis['text_shape'])

            # Collect graphic elements
            if shape_analysis['graphic_element']:
                if isinstance(shape_analysis['graphic_element'], list):
                    graphic_elements.extend(shape_analysis['graphic_element'])
                else:
                    graphic_elements.append(shape_analysis['graphic_element'])

            # Detect logos and brand elements
            if shape_analysis['logo_brand']:
                if isinstance(shape_analysis['logo_brand'], list):
                    logos_and_brands.extend(shape_analysis['logo_brand'])
                else:
                    logos_and_brands.append(shape_analysis['logo_brand'])

        # Extract notes
        notes = None
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        # Apply tags based on content
        tags = self._generate_slide_tags(text_content, notes)

        return SlideInfo(
            slide_number=slide_number,
            title=title,
            text_content=text_content,
            shape_count=shape_count,
            image_count=image_count,
            image_files=image_files,
            text_shapes=text_shapes,
            graphic_elements=graphic_elements,
            logos_and_brands=logos_and_brands,
            tags=tags,
            notes=notes
        )

    def _extract_slide_title(self, slide) -> Optional[str]:
        """Extract slide title if available."""
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
        return None

    def _extract_image(self, shape, slide_number: int) -> Optional[str]:
        """Extract and save image from shape."""
        try:
            if hasattr(shape, 'image'):
                image = shape.image
                image_bytes = image.blob

                # Create unique filename
                image_hash = hashlib.md5(image_bytes).hexdigest()[:8]
                ext = self._get_image_extension(image.content_type)
                filename = f"slide_{slide_number:02d}_{image_hash}.{ext}"

                # Save image
                image_path = self.image_dir / filename
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)

                return filename
        except Exception as e:
            print(f"Error extracting image from slide {slide_number}: {e}")

        return None

    def _analyze_text_shape(self, shape, text: str, slide_number: int) -> Optional[Dict[str, Any]]:
        """Analyze text shape formatting and properties."""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return None

            text_frame = shape.text_frame

            # Get font and formatting info
            font_info = {}
            if text_frame.paragraphs:
                first_para = text_frame.paragraphs[0]
                if first_para.runs:
                    first_run = first_para.runs[0]
                    font = first_run.font
                    font_info = {
                        'name': font.name,
                        'size': font.size.pt if font.size else None,
                        'bold': font.bold,
                        'italic': font.italic,
                        'color': str(font.color.rgb) if font.color and hasattr(font.color, 'rgb') else None
                    }

            return {
                'text': text,
                'font_info': font_info,
                'positioning': self._get_shape_positioning(shape),
                'shape_type': 'text_shape',
                'slide_number': slide_number
            }
        except Exception as e:
            return None

    def _is_logo_or_brand_text(self, text: str, shape) -> bool:
        """Detect if text might be a logo or brand element."""
        text_lower = text.lower()

        # Common logo/brand indicators
        brand_indicators = [
            'garden', 'giving', 'company', 'corp', 'inc', 'llc', 'ltd',
            '©', 'copyright', 'trademark', '™', '®', 'all rights reserved'
        ]

        # Check if text contains brand indicators
        has_brand_text = any(indicator in text_lower for indicator in brand_indicators)

        # Check positioning (logos often at top/bottom of slide)
        try:
            if hasattr(shape, 'top') and hasattr(shape, 'left'):
                # Consider it a potential logo if positioned like one
                is_positioned_like_logo = (
                    shape.top.inches < 1.0 or  # Top of slide
                    shape.top.inches > 6.0     # Bottom of slide
                )
                return has_brand_text or is_positioned_like_logo
        except:
            pass

        return has_brand_text

    def _is_logo_image(self, shape) -> bool:
        """Detect if image might be a logo."""
        try:
            # Logos are typically smaller images
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                width_inches = shape.width.inches if shape.width else 0
                height_inches = shape.height.inches if shape.height else 0

                # Consider small to medium sized images as potential logos
                is_logo_size = (width_inches < 4 and height_inches < 4) and (width_inches > 0.5 or height_inches > 0.5)

                # Check positioning
                if hasattr(shape, 'top'):
                    top_inches = shape.top.inches if shape.top else 0
                    is_logo_position = top_inches < 2.0 or top_inches > 6.0
                    return is_logo_size and is_logo_position

                return is_logo_size
        except:
            pass
        return False

    def _get_shape_positioning(self, shape) -> Dict[str, Any]:
        """Get shape positioning information."""
        try:
            return {
                'left': shape.left.inches if hasattr(shape, 'left') and shape.left else None,
                'top': shape.top.inches if hasattr(shape, 'top') and shape.top else None,
                'width': shape.width.inches if hasattr(shape, 'width') and shape.width else None,
                'height': shape.height.inches if hasattr(shape, 'height') and shape.height else None
            }
        except:
            return {}

    def _analyze_graphic_element(self, shape, slide_number: int) -> Optional[Dict[str, Any]]:
        """Analyze graphic elements like shapes and drawings."""
        try:
            element_info = {
                'type': 'graphic_element',
                'shape_type': str(shape.shape_type),
                'positioning': self._get_shape_positioning(shape),
                'slide_number': slide_number
            }

            # Try to get additional shape properties
            if hasattr(shape, 'auto_shape_type'):
                element_info['auto_shape_type'] = str(shape.auto_shape_type)

            # Check for fill color
            if hasattr(shape, 'fill'):
                try:
                    if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                        element_info['fill_color'] = str(shape.fill.fore_color.rgb)
                except:
                    pass

            # Check if it has text
            if hasattr(shape, 'text') and shape.text:
                element_info['text'] = shape.text.strip()

            return element_info
        except Exception as e:
            return None

    def _analyze_shape_comprehensive(self, shape, slide_number: int) -> Dict[str, Any]:
        """Comprehensive shape analysis for all content types."""
        result = {
            'text': None,
            'image_file': None,
            'text_shape': None,
            'graphic_element': None,
            'logo_brand': None
        }

        # Handle text content
        if hasattr(shape, 'text') and shape.text:
            text = shape.text.strip()
            result['text'] = text

            # Analyze text formatting and detect logos/brands
            text_analysis = self._analyze_text_shape(shape, text, slide_number)
            if text_analysis:
                result['text_shape'] = text_analysis

                # Detect if this might be a logo or brand element
                if self._is_logo_or_brand_text(text, shape):
                    result['logo_brand'] = {
                        'type': 'text_logo',
                        'content': text,
                        'font_info': text_analysis.get('font_info', {}),
                        'positioning': text_analysis.get('positioning', {}),
                        'slide_number': slide_number
                    }

        # Handle pictures
        if (MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.PICTURE) or shape.shape_type == 13:
            image_file = self._extract_image(shape, slide_number)
            if image_file:
                result['image_file'] = image_file

                # Check if image might be a logo
                if self._is_logo_image(shape):
                    result['logo_brand'] = {
                        'type': 'image_logo',
                        'file': image_file,
                        'positioning': self._get_shape_positioning(shape),
                        'slide_number': slide_number
                    }

        # Handle other graphic elements (shapes, drawings, etc.)
        elif (MSO_SHAPE_TYPE and shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]) or shape.shape_type in [1, 5]:
            graphic_info = self._analyze_graphic_element(shape, slide_number)
            if graphic_info:
                result['graphic_element'] = graphic_info

        # Handle group shapes recursively
        elif (MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.GROUP) or shape.shape_type == 6:
            group_results = self._process_group_shape_comprehensive(shape, slide_number)
            # Merge group results
            for key in result.keys():
                if group_results.get(key):
                    if isinstance(result[key], list):
                        result[key].extend(group_results[key])
                    elif result[key] is None:
                        result[key] = group_results[key]

        return result

    def _process_group_shape_comprehensive(self, group_shape, slide_number: int) -> Dict[str, List]:
        """Process all content within group shapes."""
        results = {
            'text': [],
            'image_file': [],
            'text_shape': [],
            'graphic_element': [],
            'logo_brand': []
        }

        for shape in group_shape.shapes:
            shape_analysis = self._analyze_shape_comprehensive(shape, slide_number)
            for key in results.keys():
                if shape_analysis[key]:
                    if isinstance(shape_analysis[key], list):
                        results[key].extend(shape_analysis[key])
                    else:
                        results[key].append(shape_analysis[key])

        return results

    def _get_image_extension(self, content_type: str) -> str:
        """Get file extension from content type."""
        type_map = {
            'image/png': 'png',
            'image/jpeg': 'jpg',
            'image/gif': 'gif',
            'image/bmp': 'bmp',
            'image/tiff': 'tiff'
        }
        return type_map.get(content_type, 'png')

    def _get_all_text_content(self) -> str:
        """Get all text content from presentation."""
        all_text = []

        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    all_text.append(shape.text)

        return ' '.join(all_text).lower()

    def _count_total_images(self) -> int:
        """Count total images in presentation."""
        total = 0
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    total += 1
                elif hasattr(shape, 'shapes'):  # Group
                    total += self._count_images_in_group(shape)
        return total

    def _count_images_in_group(self, group_shape) -> int:
        """Count images in group shape."""
        count = 0
        for shape in group_shape.shapes:
            if shape.shape_type == 13:
                count += 1
            elif hasattr(shape, 'shapes'):
                count += self._count_images_in_group(shape)
        return count

    def _find_pattern_matches(self, text: str, patterns: List[str]) -> List[str]:
        """Find pattern matches in text."""
        matches = []
        for pattern in patterns:
            if pattern.lower() in text.lower():
                matches.append(pattern)
        return list(set(matches))  # Remove duplicates

    def _generate_slide_tags(self, text_content: List[str], notes: Optional[str]) -> List[str]:
        """Generate tags for slide based on content."""
        tags = []
        all_text = ' '.join(text_content + ([notes] if notes else [])).lower()

        # Check for copyright
        if any(pattern in all_text for pattern in self.copyright_patterns):
            tags.append("copyright")

        # Check for confidentiality
        if any(pattern in all_text for pattern in self.confidentiality_patterns):
            tags.append("confidential")

        # Check for company mentions
        if any(pattern in all_text for pattern in self.company_patterns):
            tags.append("company_info")

        # Add content-based tags
        if "agenda" in all_text:
            tags.append("agenda")
        if any(word in all_text for word in ["introduction", "overview"]):
            tags.append("introduction")
        if any(word in all_text for word in ["conclusion", "summary", "thank you"]):
            tags.append("conclusion")

        return tags

    def apply_standardized_tags(self, custom_tags: Dict[str, List[str]] = None) -> None:
        """Apply standardized tags across presentation."""
        if custom_tags:
            self.metadata.custom_tags.extend(custom_tags.get("document", []))

            for i, slide_info in enumerate(self.slides_info):
                slide_tags = custom_tags.get(f"slide_{slide_info.slide_number}", [])
                slide_info.tags.extend(slide_tags)

    def export_to_json(self, filename: str = None) -> str:
        """Export all data to JSON format."""
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{self.filepath.stem}_analysis_{timestamp}.json"

        output_path = self.output_dir / filename

        data = {
            "metadata": asdict(self.metadata),
            "slides": [asdict(slide) for slide in self.slides_info],
            "analysis_timestamp": datetime.now().isoformat(),
            "inspector_version": "1.0.0"
        }

        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2, ensure_ascii=False)

        print(f"JSON export saved: {output_path}")
        return str(output_path)

    def export_to_csv(self, filename: str = None) -> str:
        """Export slide data to CSV format."""
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{self.filepath.stem}_slides_{timestamp}.csv"

        output_path = self.output_dir / filename

        fieldnames = [
            "slide_number",
            "title",
            "text_content",
            "shape_count",
            "image_count",
            "image_files",
            "tags",
            "notes"
        ]

        with open(output_path, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
            writer.writeheader()

            for slide in self.slides_info:
                writer.writerow({
                    "slide_number": slide.slide_number,
                    "title": slide.title or "",
                    "text_content": " | ".join(slide.text_content),
                    "shape_count": slide.shape_count,
                    "image_count": slide.image_count,
                    "image_files": " | ".join(slide.image_files),
                    "tags": " | ".join(slide.tags),
                    "notes": slide.notes or ""
                })

        print(f"CSV export saved: {output_path}")
        return str(output_path)

    def generate_summary_report(self) -> Dict[str, Any]:
        """Generate summary report."""
        if not self.metadata or not self.slides_info:
            raise ValueError("No data available. Run analysis first.")

        # Aggregate statistics
        total_text_items = sum(len(slide.text_content) for slide in self.slides_info)
        slides_with_images = sum(1 for slide in self.slides_info if slide.image_count > 0)
        all_tags = set()
        for slide in self.slides_info:
            all_tags.update(slide.tags)

        summary = {
            "file_info": {
                "name": self.metadata.filename,
                "size_mb": round(self.metadata.file_size / (1024 * 1024), 2),
                "slide_count": self.metadata.slide_count
            },
            "content_analysis": {
                "total_text_items": total_text_items,
                "total_images": self.metadata.total_images,
                "slides_with_images": slides_with_images,
                "unique_tags": list(all_tags)
            },
            "compliance_check": {
                "has_copyright": len(self.metadata.copyright_notices) > 0,
                "has_confidentiality": len(self.metadata.confidentiality_labels) > 0,
                "company_mentions": self.metadata.company_mentions
            }
        }

        return summary


@click.command()
@click.argument('filepath', type=click.Path(exists=True))
@click.option('--output-dir', '-o', default='exports', help='Output directory for exports')
@click.option('--image-dir', '-i', default='images', help='Directory for extracted images')
@click.option('--tags-file', '-t', type=click.Path(exists=True), help='JSON file with custom tags')
@click.option('--export-format', '-f', multiple=True, default=['json'],
              type=click.Choice(['json', 'csv', 'both']), help='Export format(s)')
def main(filepath, output_dir, image_dir, tags_file, export_format):
    """PowerPoint Inspector - Extract and analyze PowerPoint presentations."""

    print(f"🔍 Analyzing PowerPoint file: {filepath}")

    # Initialize inspector
    inspector = PowerPointInspector(filepath, output_dir, image_dir)

    # Load and analyze presentation
    inspector.load_presentation()

    print("\n📊 Extracting document metadata...")
    metadata = inspector.extract_document_metadata()

    print("📄 Analyzing slide content...")
    slides_info = inspector.extract_slide_content()

    # Apply custom tags if provided
    if tags_file:
        print(f"🏷️  Applying custom tags from: {tags_file}")
        with open(tags_file, 'r') as f:
            custom_tags = json.load(f)
        inspector.apply_standardized_tags(custom_tags)

    # Generate summary
    print("\n📈 Generating summary report...")
    summary = inspector.generate_summary_report()

    print("\n=== SUMMARY REPORT ===")
    print(f"File: {summary['file_info']['name']} ({summary['file_info']['size_mb']} MB)")
    print(f"Slides: {summary['file_info']['slide_count']}")
    print(f"Text Items: {summary['content_analysis']['total_text_items']}")
    print(f"Images: {summary['content_analysis']['total_images']} (in {summary['content_analysis']['slides_with_images']} slides)")
    print(f"Tags Found: {', '.join(summary['content_analysis']['unique_tags']) if summary['content_analysis']['unique_tags'] else 'None'}")
    print(f"Copyright: {'✓' if summary['compliance_check']['has_copyright'] else '✗'}")
    print(f"Confidentiality: {'✓' if summary['compliance_check']['has_confidentiality'] else '✗'}")

    # Export data
    print("\n💾 Exporting data...")
    for fmt in export_format:
        if fmt == 'json' or fmt == 'both':
            inspector.export_to_json()
        if fmt == 'csv' or fmt == 'both':
            inspector.export_to_csv()

    print("\n✅ Analysis complete!")


if __name__ == "__main__":
    main()
